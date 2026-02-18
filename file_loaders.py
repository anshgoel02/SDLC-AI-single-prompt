from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pypdf import PdfReader

from .config import SUPPORTED_EXTENSIONS

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from docx import Document
except Exception:
    Document = None


# Extracts and concatenates text from all pages in a PDF file.
def read_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


# Extracts text content from all shapes across all slides in a PPTX file.
def read_pptx_text(path: Path) -> str:
    if Presentation is None:
        return ""
    pres = Presentation(str(path))
    parts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


# Extracts paragraph text from a DOCX file.
def read_docx_text(path: Path) -> str:
    if Document is None:
        return ""
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


# Loads text content from supported non-image files by extension.
def load_source_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf_text(path)
    if suffix == ".pptx":
        return read_pptx_text(path)
    if suffix == ".docx":
        return read_docx_text(path)
    return path.read_text(encoding="utf-8", errors="ignore")


# Returns whether the path points to a supported image format.
def is_image_file(path: Path) -> bool:
    return path.suffix.lower() in {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".gif",
        ".webp",
        ".tif",
        ".tiff",
    }


# Yields supported files from input paths (files or recursive directories).
def iter_input_files(paths: Sequence[str]) -> Iterable[Path]:
    for raw in paths:
        path = Path(raw)
        if not path.exists():
            continue
        if path.is_file():
            if path.suffix.lower() in SUPPORTED_EXTENSIONS:
                yield path
        else:
            for file_path in path.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                    yield file_path
